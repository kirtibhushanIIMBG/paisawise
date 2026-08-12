#!/usr/bin/env python3
"""Build PaisaWise_Pitch_Deck_v2.pptx directly with python-pptx.

5 slides, 3 pitchers, 5 minutes. No HTML intermediate.

Colour system: indigo-ink + electric violet, mint for positives.
Slide 1 is the dark hero cover; slides 2-5 share one light surface so the
body of the deck is visually uniform.

Run:  python3 build_deck.py
"""
import os
import cairosvg
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE  = os.path.dirname(os.path.abspath(__file__))
BRAND = os.path.join(HERE, "brand")
ICONS = os.path.join(BRAND, "icons")
OUT   = os.path.join(HERE, "PaisaWise_Pitch_Deck_v2.pptx")
os.makedirs(ICONS, exist_ok=True)

# ----------------------------------------------------------------- palette
INK      = "171526"   # near-black indigo — cover, bands, headlines
INK_2    = "241F3D"   # raised surface on ink
INK_3    = "3A3163"   # border / divider on ink
SURFACE  = "F6F5FA"   # light slide ground (slides 2-5)
CARD     = "FFFFFF"
CARD_ALT = "FBFAFD"
BORDER   = "E6E3F0"
BODY     = "5A5474"   # body text on light
MUTED    = "9A94B8"   # secondary text
PALE     = "C9C3E4"   # body text on ink
VIOLET   = "6C4CF1"   # primary accent
VIOLET_L = "9B85FF"   # accent on ink
VIOLET_T = "EFEAFF"   # accent tint
VIOLET_T2= "F5F2FF"   # accent tint, alternate row
MINT     = "12B981"   # positive / gain
WHITE    = "FFFFFF"

F = "Aptos"
W, H = 13.333, 7.5
M = 0.60
CW = W - 2 * M

def rgb(h):
    return RGBColor.from_string(h)

# ----------------------------------------------------------------- icons
ICON_SVG = {
 "eye":     '<path d="M2 12s3.6-6.5 10-6.5S22 12 22 12s-3.6 6.5-10 6.5S2 12 2 12Z"/><circle cx="12" cy="12" r="2.8"/>',
 "chat":    '<path d="M21 14.5a2.2 2.2 0 0 1-2.2 2.2H8.2L4 20.5V5.2A2.2 2.2 0 0 1 6.2 3h12.6A2.2 2.2 0 0 1 21 5.2Z"/>',
 "user":    '<circle cx="12" cy="8" r="3.6"/><path d="M4.8 20.5v-1a5.2 5.2 0 0 1 5.2-5.2h4a5.2 5.2 0 0 1 5.2 5.2v1"/>',
 "grid":    '<rect x="3" y="3" width="7.5" height="7.5" rx="1.6"/><rect x="13.5" y="3" width="7.5" height="7.5" rx="1.6"/><rect x="3" y="13.5" width="7.5" height="7.5" rx="1.6"/><rect x="13.5" y="13.5" width="7.5" height="7.5" rx="1.6"/>',
 "chip":    '<rect x="6.5" y="6.5" width="11" height="11" rx="2.4"/><path d="M9.5 3v3.5M14.5 3v3.5M9.5 17.5V21M14.5 17.5V21M3 9.5h3.5M3 14.5h3.5M17.5 9.5H21M17.5 14.5H21"/>',
 "target":  '<circle cx="12" cy="12" r="8.6"/><circle cx="12" cy="12" r="4.6"/><circle cx="12" cy="12" r="1.3" fill="CUR" stroke="none"/>',
 "bell":    '<path d="M18 8.6a6 6 0 1 0-12 0c0 6.4-2.6 8.4-2.6 8.4h17.2S18 15 18 8.6"/><path d="M13.8 20.6a2.1 2.1 0 0 1-3.6 0"/>',
 "pie":     '<circle cx="12" cy="12" r="8.8"/><path d="M12 3.2V12l7.6 4.4"/>',
 "advisor": '<circle cx="10.5" cy="8" r="3.4"/><path d="M3.8 20.4v-.9a5 5 0 0 1 5-5h3.6"/><path d="m15.2 17.6 2 2 4-4"/>',
 "check":   '<path d="m20 6.5-10.5 11L4 12"/>',
 "cross":   '<path d="M17.5 6.5l-11 11M6.5 6.5l11 11"/>',
 "bank":    '<path d="M12 3.2 21 8H3ZM4.8 8v9M9.6 8v9M14.4 8v9M19.2 8v9M2.6 20.8h18.8"/>',
 "store":   '<path d="M3.4 9.2 5 4h14l1.6 5.2M4.6 9.2v11.2h14.8V9.2M9.6 20.4v-5.6h4.8v5.6"/>',
 "linkedin":'<rect x="3" y="3" width="18" height="18" rx="4"/><circle cx="7.6" cy="8" r="1.15" fill="CUR" stroke="none"/><path d="M7.6 11v6.5M11.7 17.5V11M11.7 14.1a2.6 2.6 0 0 1 5.2 0v3.4"/>',
 "users":   '<circle cx="9" cy="8" r="3.2"/><path d="M2.8 20.4v-.8a4.8 4.8 0 0 1 4.8-4.8h2.8a4.8 4.8 0 0 1 4.8 4.8v.8"/><path d="M16.4 5.2a3.2 3.2 0 0 1 0 5.9M18 14.9a4.8 4.8 0 0 1 3.2 4.5v1"/>',
 "briefcase":'<rect x="2.8" y="7.2" width="18.4" height="12" rx="2.2"/><path d="M8.8 7.2V5.4a2 2 0 0 1 2-2h2.4a2 2 0 0 1 2 2v1.8"/><path d="M2.8 12.6h18.4"/>',
 "trend":   '<path d="M3.5 16.5 9 11l3.5 3.5L20.5 6.5"/><path d="M15.5 6.5h5v5"/>',
 "clock":   '<circle cx="12" cy="12" r="8.8"/><path d="M12 6.8V12l3.6 2.2"/>',
}

def icon(name, colour, px=180):
    path = os.path.join(ICONS, f"{name}-{colour}.png")
    if not os.path.exists(path):
        body = ICON_SVG[name].replace("CUR", f"#{colour}")
        svg = (f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 24 24" '
               f'width="24" height="24" fill="none" stroke="#{colour}" stroke-width="1.9" '
               f'stroke-linecap="round" stroke-linejoin="round">{body}</svg>')
        cairosvg.svg2png(bytestring=svg.encode(), write_to=path,
                         output_width=px, output_height=px)
    return path

# ----------------------------------------------------------------- primitives
def bg(slide, colour):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(W), Inches(H))
    s.fill.solid(); s.fill.fore_color.rgb = rgb(colour); s.line.fill.background()
    s.shadow.inherit = False
    return s

def rect(slide, x, y, w, h, fill=None, line=None, radius=0.0, lw=1.0):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if radius:
        s.adjustments[0] = min(0.5, radius / (min(w, h) / 2.0))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = rgb(fill)
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = rgb(line); s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s

def ellipse(slide, x, y, d, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    s.fill.solid(); s.fill.fore_color.rgb = rgb(fill); s.line.fill.background()
    s.shadow.inherit = False
    return s

def txt(slide, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p.get("align", align)
        if p.get("line"):   para.line_spacing = p["line"]
        if p.get("before"): para.space_before = Pt(p["before"])
        if p.get("after"):  para.space_after = Pt(p["after"])
        for text, st in p["runs"]:
            r = para.add_run(); r.text = text
            f = r.font
            f.name = F
            f.size = Pt(st.get("size", 12))
            f.bold = st.get("bold", False)
            f.italic = st.get("italic", False)
            f.color.rgb = rgb(st.get("color", INK))
            if st.get("spc"):
                f._element.set("spc", str(int(st["spc"] * 100)))
    return box

def one(slide, x, y, w, h, text, size=12, color=INK, bold=False, italic=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spc=None, line=None):
    return txt(slide, x, y, w, h,
               [{"runs": [(text, dict(size=size, color=color, bold=bold,
                                      italic=italic, spc=spc))], "line": line}],
               align, anchor)

def eyebrow(slide, text, colour=VIOLET):
    one(slide, M, 0.40, 8.0, 0.30, text, size=11, color=colour, bold=True, spc=1.6)

def pic(slide, path, x, y, w=None, h=None):
    kw = {}
    if w: kw["width"] = Inches(w)
    if h: kw["height"] = Inches(h)
    return slide.shapes.add_picture(path, Inches(x), Inches(y), **kw)

def corner_mark(slide, dark=False):
    f = "paisawise-mark-reverse.png" if dark else "paisawise-mark.png"
    pic(slide, os.path.join(BRAND, f), W - M - 0.46, 0.34, w=0.46)

def badge_icon(slide, name, x, y, d=0.50, dark=False):
    """Tinted disc with a line icon centred inside it."""
    ellipse(slide, x, y, d, INK_2 if dark else VIOLET_T)
    pic(slide, icon(name, VIOLET_L if dark else VIOLET), x + d / 4, y + d / 4, w=d / 2)

def chip(slide, x, y, w, h, label, dark=False, size=10.5, ic=None):
    fill = INK_2 if dark else CARD
    line = INK_3 if dark else BORDER
    col  = PALE if dark else INK
    rect(slide, x, y, w, h, fill, line, radius=h / 2)
    tx, tw = x, w
    if ic:
        pic(slide, icon(ic, VIOLET_L if dark else VIOLET), x + 0.30, y + (h - 0.24) / 2, w=0.24)
        tx, tw = x + 0.66, w - 0.90
    one(slide, tx, y, tw, h, label, size=size, color=col, bold=True,
        align=PP_ALIGN.LEFT if ic else PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def label_band(slide, y, h, label, text, size=15, italic=False):
    """Ink band with a violet kicker label, then the line itself."""
    rect(slide, M, y, CW, h, INK, radius=0.12)
    txt(slide, M + 0.40, y, CW - 0.80, h,
        [{"runs": [(label + "      ", dict(size=11, color=VIOLET_L, bold=True, spc=1.6)),
                   (text, dict(size=size, color=WHITE, bold=True, italic=italic))]}],
        anchor=MSO_ANCHOR.MIDDLE)

def bullets(slide, x, y, w, items, dark=False, size=11.5, step=0.42):
    dot = VIOLET_L if dark else VIOLET
    col = PALE if dark else BODY
    for i, t in enumerate(items):
        txt(slide, x, y + i * step, w, step,
            [{"runs": [("•  ", dict(size=size, color=dot, bold=True)),
                       (t, dict(size=size, color=col))]}],
            anchor=MSO_ANCHOR.MIDDLE)

# ================================================================= slides
prs = Presentation()
prs.slide_width, prs.slide_height = Inches(W), Inches(H)
BLANK = prs.slide_layouts[6]
def new():
    return prs.slides.add_slide(BLANK)

# ----------------------------------------------------------------- 1 · cover
def slide1():
    s = new(); bg(s, INK)
    ellipse(s, -1.9, 4.4, 4.6, INK_2)
    ellipse(s, 11.2, -1.5, 4.2, INK_2)

    logo = os.path.join(BRAND, "paisawise-logo-horizontal-reverse.png")
    from PIL import Image
    iw, ih = Image.open(logo).size
    lh = 0.92; lw = lh * iw / ih
    pic(s, logo, (W - lw) / 2, 0.72, h=lh)

    txt(s, 0, 2.24, W, 0.60,
        [{"runs": [("You earn well.", dict(size=38, color=WHITE, bold=True))]}],
        PP_ALIGN.CENTER)
    txt(s, 0, 2.88, W, 0.72,
        [{"runs": [("So where does it ", dict(size=38, color=WHITE, bold=True)),
                   ("all go?", dict(size=38, color=VIOLET_L, bold=True))]}],
        PP_ALIGN.CENTER)

    one(s, 2.40, 3.94, 8.53, 0.72,
        "An AI budgeting app with a certified financial advisor attached. Built for "
        "salaried professionals who earn well and still can't answer that question.",
        size=13.5, color=PALE, align=PP_ALIGN.CENTER, line=1.35)

    cw, gap = 2.70, 0.28
    x0 = (W - (3 * cw + 2 * gap)) / 2
    for i, lbl in enumerate(["Syncs your banks + UPI",
                             "Certified advisor monthly",
                             "₹499 a month"]):
        chip(s, x0 + i * (cw + gap), 5.06, cw, 0.62, lbl, dark=True, size=12)

    rect(s, (W - 4.2) / 2, 6.06, 4.2, 0.022, VIOLET)
    agenda = [("PITCHER 1", "The Problem"), ("PITCHER 2", "The Solution"),
              ("PITCHER 3", "Why PaisaWise & The Close")]
    runs = []
    for i, (p, t) in enumerate(agenda):
        if i: runs.append(("     •     ", dict(size=10.5, color=INK_3, bold=True)))
        runs += [(p + " ", dict(size=10.5, color=VIOLET_L, bold=True, spc=0.8)),
                 (t, dict(size=10.5, color=MUTED))]
    txt(s, 0, 6.36, W, 0.34, [{"runs": runs}], PP_ALIGN.CENTER)
    return s

# ----------------------------------------------------------------- 2 · problem
def slide2():
    s = new(); bg(s, SURFACE)
    eyebrow(s, "PITCHER 1  ·  THE PROBLEM"); corner_mark(s)

    one(s, M, 0.80, 11.4, 0.48, "Two problems. Everything else follows from them.",
        size=26, color=INK, bold=True)
    one(s, M, 1.40, 11.4, 0.34,
        "Salaried professional, 25 to 35, introduced by his bank relationship manager after a branch visit.",
        size=12.5, color=BODY)

    px, pw, py, ph = M, 4.05, 1.92, 3.58
    rect(s, px, py, pw, ph, CARD, BORDER, radius=0.12)
    badge_icon(s, "user", px + 0.30, py + 0.28)
    one(s, px + 0.96, py + 0.28, pw - 1.2, 0.50, "THE BUYER",
        size=11, color=VIOLET, bold=True, spc=1.4, anchor=MSO_ANCHOR.MIDDLE)
    bullets(s, px + 0.30, py + 0.98, pw - 0.60, [
        "Age 25–35, salaried",
        "Middle to upper-middle income",
        "Regular bank + UPI transactions",
        "Earns well, struggles to budget",
        "Knows SIPs, lacks confidence",
    ], size=11.5, step=0.47)

    cx, cwd = 4.95, 7.78
    pains = [
        ("eye", "Poor visibility of spending",
         "His money moves across three or four bank accounts and a couple of UPI apps. "
         "Nothing adds it up for him."),
        ("chat", "Generic financial advice",
         "Free apps give the same tips to everyone who installs them. None of it knows "
         "how he earns or where he overspends."),
    ]
    for i, (ic, title, body) in enumerate(pains):
        y = 1.92 + i * 1.88
        rect(s, cx, y, cwd, 1.70, CARD, BORDER, radius=0.12)
        rect(s, cx, y + 0.22, 0.055, 1.26, VIOLET)
        badge_icon(s, ic, cx + 0.34, y + 0.26)
        one(s, cx + 1.00, y + 0.26, cwd - 1.3, 0.50, title,
            size=15, color=INK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        one(s, cx + 0.34, y + 0.88, cwd - 0.70, 0.64, body,
            size=11.5, color=BODY, line=1.30)

    one(s, M, 5.66, 5.0, 0.26, "ALSO COMES UP", size=9.5, color=MUTED, bold=True, spc=1.2)
    secondary = ["Struggles to save", "Low investing confidence",
                 "No time to plan", "Alerts arrive too late"]
    cwid = (CW - 3 * 0.23) / 4
    for i, t in enumerate(secondary):
        chip(s, M + i * (cwid + 0.23), 5.98, cwid, 0.44, t, size=10.5)

    label_band(s, 6.62, 0.62, "WHAT WE HEAR ON THE FIRST CALL",
               "“I know I should be investing. I just don't know where to start.”",
               size=14.5, italic=True)
    return s

# ----------------------------------------------------------------- 3 · solution
def slide3():
    s = new(); bg(s, SURFACE)
    eyebrow(s, "PITCHER 2  ·  THE SOLUTION"); corner_mark(s)

    one(s, M, 0.78, 11.4, 0.48, "What he gets for ₹499 a month",
        size=26, color=INK, bold=True)

    feats = [
        ("grid",    "One dashboard",             "All his accounts and UPI apps in a single view."),
        ("chip",    "AI-driven budgeting",       "A budget built from his own spending history, not a template."),
        ("target",  "Savings goals that fit",    "Monthly targets set from what he earns and spends."),
        ("bell",    "Overspending alerts",       "He hears about it the same day, not at month end."),
        ("pie",     "SIP and mutual fund nudges","Matched to his risk profile."),
        ("advisor", "A certified advisor, monthly", "One-on-one video call to go through his numbers."),
    ]
    cwid = (CW - 2 * 0.25) / 3
    for i, (ic, title, body) in enumerate(feats):
        col, row = i % 3, i // 3
        x = M + col * (cwid + 0.25)
        y = 1.52 + row * 1.98
        rect(s, x, y, cwid, 1.82, CARD, BORDER, radius=0.12)
        badge_icon(s, ic, x + 0.28, y + 0.28)
        one(s, x + 0.90, y + 0.28, cwid - 1.15, 0.50, title,
            size=13, color=INK, bold=True, anchor=MSO_ANCHOR.MIDDLE, line=1.05)
        one(s, x + 0.28, y + 0.94, cwid - 0.56, 0.72, body,
            size=10.5, color=BODY, line=1.32)

    label_band(s, 5.96, 0.88, "OUR USP",
               "AI + personalization + a human financial advisor", size=16)
    return s

# ----------------------------------------------------------------- 4 · why
def slide4():
    s = new(); bg(s, SURFACE)
    eyebrow(s, "PITCHER 3  ·  WHY PAISAWISE"); corner_mark(s)

    one(s, M, 0.78, 11.4, 0.48, "Why not just use a free app?",
        size=26, color=INK, bold=True)
    one(s, M, 1.36, 11.4, 0.34, "He already has one on his phone.",
        size=12.5, color=BODY)

    tx, tw = M, 7.10
    half = tw / 2
    ty, hh, rh = 1.90, 0.52, 0.72
    rect(s, tx, ty, half, hh, BORDER)
    rect(s, tx + half, ty, half, hh, INK)
    one(s, tx, ty, half, hh, "FREE TRACKING APPS", size=10.5, color=BODY, bold=True,
        spc=1.2, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    one(s, tx + half, ty, half, hh, "PAISAWISE", size=10.5, color=WHITE, bold=True,
        spc=1.2, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    rows = [
        ("Shows where money went",     "Tells you what to do next"),
        ("Generic insights",           "Personalized to your spending"),
        ("No live coaching",           "Monthly 1-on-1 advisor call"),
        ("No proactive alerts",        "Real-time overspending alerts"),
        ("Limited investment guidance","Risk-matched SIP / MF nudges"),
    ]
    for i, (l, r) in enumerate(rows):
        y = ty + hh + i * rh
        rect(s, tx, y, half, rh, CARD if i % 2 == 0 else CARD_ALT, BORDER)
        rect(s, tx + half, y, half, rh, VIOLET_T if i % 2 == 0 else VIOLET_T2, BORDER)
        pic(s, icon("cross", MUTED), tx + 0.22, y + (rh - 0.20) / 2, w=0.20)
        one(s, tx + 0.54, y, half - 0.70, rh, l, size=11, color=BODY,
            anchor=MSO_ANCHOR.MIDDLE)
        pic(s, icon("check", MINT), tx + half + 0.22, y + (rh - 0.20) / 2, w=0.20)
        one(s, tx + half + 0.54, y, half - 0.70, rh, r, size=11, color=INK, bold=True,
            anchor=MSO_ANCHOR.MIDDLE)

    bx, bw, by, bh = 7.95, 4.78, 1.90, 4.12
    rect(s, bx, by, bw, bh, INK, radius=0.12)
    one(s, bx + 0.34, by + 0.34, bw - 0.68, 0.26, "THE MATH",
        size=10.5, color=VIOLET_L, bold=True, spc=1.4)
    txt(s, bx + 0.34, by + 0.86, bw - 0.68, 0.48,
        [{"runs": [("₹499", dict(size=30, color=WHITE, bold=True)),
                   ("/month", dict(size=13, color=MUTED)),
                   ("   ≈  ₹16/day", dict(size=13, color=PALE, bold=True))]}],
        anchor=MSO_ANCHOR.MIDDLE)
    rect(s, bx + 0.34, by + 1.58, bw - 0.68, 0.012, INK_3)
    one(s, bx + 0.34, by + 1.80, bw - 0.68, 0.46,
        "If it helps him save 5% of a ₹60,000 salary",
        size=11.5, color=PALE, line=1.28)
    txt(s, bx + 0.34, by + 2.44, bw - 0.68, 0.44,
        [{"runs": [("₹3,000", dict(size=26, color=MINT, bold=True)),
                   (" back, every month", dict(size=13, color=PALE))]}],
        anchor=MSO_ANCHOR.MIDDLE)
    one(s, bx + 0.34, by + 3.02, bw - 0.68, 0.28, "Six times what it costs him.",
        size=13, color=WHITE, bold=True)
    one(s, bx + 0.34, by + 3.62, bw - 0.68, 0.26,
        "Illustration, not a guarantee.", size=9, color=MUTED, italic=True)

    one(s, M, 6.42, 9.8, 0.34,
        "Free apps are ad-supported: they earn when he scrolls. We only earn when he renews.",
        size=13, color=BODY)
    return s

# ----------------------------------------------------------------- 5 · close
def slide5():
    s = new(); bg(s, SURFACE)
    eyebrow(s, "PITCHER 3  ·  THE CLOSE"); corner_mark(s)

    one(s, M, 0.78, 11.4, 0.48, "Who we sell to, and how we reach them",
        size=26, color=INK, bold=True)

    cwd = (CW - 0.33) / 2
    rect(s, M, 1.48, cwd, 2.92, CARD, BORDER, radius=0.12)
    one(s, M + 0.34, 1.70, cwd - 0.68, 0.28, "THE IDEAL LEAD",
        size=11, color=VIOLET, bold=True, spc=1.4)
    bullets(s, M + 0.34, 2.08, cwd - 0.68, [
        "Stable, predictable income",
        "Multiple bank + UPI transactions",
        "Actively wants to save or invest",
        "Lacks personalized guidance",
        "Willing to pay for expert help",
    ], size=11.5, step=0.43)

    sx = M + cwd + 0.33
    rect(s, sx, 1.48, cwd, 2.92, CARD, BORDER, radius=0.12)
    one(s, sx + 0.34, 1.70, cwd - 0.68, 0.28, "WHERE WE FIND THEM",
        size=11, color=VIOLET, bold=True, spc=1.4)
    srcs = [("bank", "Bank relationship managers"),
            ("briefcase", "Bank & corporate partnerships"),
            ("store", "Branch walk-ins"),
            ("linkedin", "LinkedIn: working professionals"),
            ("users", "Referrals from existing customers")]
    for i, (ic, t) in enumerate(srcs):
        y = 2.08 + i * 0.43
        pic(s, icon(ic, VIOLET), sx + 0.36, y + 0.09, w=0.235)
        one(s, sx + 0.80, y, cwd - 1.14, 0.43, t, size=11.5, color=BODY,
            anchor=MSO_ANCHOR.MIDDLE)

    rect(s, M, 4.62, CW, 0.90, VIOLET, radius=0.12)
    txt(s, M + 0.42, 4.62, 5.20, 0.90,
        [{"runs": [("₹499", dict(size=27, color=WHITE, bold=True)),
                   ("/month", dict(size=13, color=VIOLET_T)),
                   ("     ₹4,999", dict(size=19, color=WHITE, bold=True)),
                   ("/year", dict(size=13, color=VIOLET_T)),
                   ("     ≈ ₹16/day", dict(size=13, color=VIOLET_T, bold=True))]}],
        anchor=MSO_ANCHOR.MIDDLE)
    one(s, M + 5.80, 4.62, CW - 6.22, 0.90,
        "₹999 onboarding fee waived for bank customers",
        size=12.5, color=WHITE, bold=True, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    rect(s, M, 5.78, CW, 1.14, INK, radius=0.12)
    rect(s, M, 5.98, 0.06, 0.74, VIOLET_L)
    one(s, M + 0.34, 5.96, 3.0, 0.26, "THE ASK", size=11, color=VIOLET_L, bold=True, spc=1.4)
    one(s, M + 0.34, 6.26, CW - 0.80, 0.52,
        "“You came in through your relationship manager, so the ₹999 onboarding fee is "
        "waived. Shall we book your first advisor call this week?”",
        size=14.5, color=WHITE, bold=True, italic=True, line=1.25)
    return s

# ----------------------------------------------------------------- notes
NOTES = [
"""PITCHER 1  ·  target 0:00–0:22  (22s)

Good morning. One question before I start. You earn well. So where did your money
go last month? Most salaried professionals cannot. That is the
gap we built PaisaWise to close. Over the next five minutes, the three of us will
show you the problem, the product, and what it's worth to you.

[CUE] Advance on "the problem".""",

"""PITCHER 1  ·  target 0:22–1:40  (78s)   ·   cumulative 1:40

Let me describe our customer, because you will recognise him. He's 25 to 35, salaried,
middle to upper-middle income. He came to us the way most of our customers do. His
bank relationship manager introduced us after a routine branch visit.

He earns well. But his money is scattered across three or four bank accounts and a
couple of UPI apps. So at the end of the month, when he asks "where did it all go?",
there is no single place that answers him. That is problem one: poor visibility.

Problem two. When he does go looking for help, he gets generic advice. Free tracking
apps give the same five tips to five million people. None of it knows that he spends
heavily in the first ten days, or that his rent went up last quarter.

Everything else follows from those two. He struggles to save. He knows about SIPs but
has never started one. He keeps deferring the planning. And he finds out he has
overspent only after the money is gone.

On most first calls what we actually hear is this: I know I should be investing, I
just don't know where to start.

[HAND-OFF] "That is what PaisaWise fixes. Over to you, <NAME 2>."
""",

"""PITCHER 2  ·  target 1:40–3:19  (99s)   ·   cumulative 3:19

Thanks. Here is what PaisaWise actually does.

First, one dashboard. We securely sync every linked bank account and UPI app, so all
his spending finally sits in a single view. That alone solves the visibility problem.

Second, AI-driven budgeting. This is the important part. The budget is not a template.
It is built from his own spending history, so it knows his patterns.

Third, personalized savings goals. Because we can see real cash flow, we set a monthly
target he can actually hit, instead of an arbitrary number he abandons in week two.

Fourth, overspending alerts. He hears about it the same day, while he can still do
something about it, instead of reading a report at month end.

Fifth, SIP and mutual fund nudges matched to his risk profile. He already knows he
should invest. We tell him what, how much, and when.

And sixth, the one no free app offers: a one-on-one virtual call every month with a
certified financial advisor, going through his numbers with him.

That is our USP in one line: AI, plus personalization, plus a human financial advisor.
The technology finds the insight. The advisor makes sure he acts on it.

[HAND-OFF] "So why us and not a free app? <NAME 3> will take that."
""",

"""PITCHER 3  ·  target 3:19–4:15  (56s)   ·   cumulative 4:15

Thanks. Now the obvious objection: there are free apps that do some of this. Let's be
precise about what free actually buys you.

A free app shows you where your money went. PaisaWise tells you what to do next. Free
gives you generic insights; we personalize to your spending. Free has no live coaching;
you get a certified advisor every month. Free has no proactive alerts; we tell you the
same day. And SIP nudges matched to your risk profile.

Now the math. To be clear, this is an illustration, not a promise. PaisaWise costs 499
rupees a month, about sixteen rupees a day. If we help you save five percent of a
sixty-thousand-rupee salary, that is three thousand rupees a month. Six times what you
paid us.

One more thing worth saying out loud: free apps are ad-supported. They make money when
you scroll. We only make money when you renew.""",

"""PITCHER 3  ·  target 4:15–5:00  (45s)   ·   cumulative 5:00   ·   TOTAL 300s

So who is the ideal lead? Someone with a stable income and plenty of bank and UPI
activity, who genuinely wants to save or invest, but has nobody personalising the
advice for him.

We reach them through bank relationship managers and corporate partnerships, branch
walk-ins, LinkedIn, and referrals from customers we have already helped.

Pricing is 499 rupees a month, or 4,999 a year. About sixteen rupees a day.

And here is my ask. You came in through your relationship manager, so the 999-rupee
onboarding fee is waived. Shall we book your first advisor call this week?

────────────────────────────────────────────────────────────
OBJECTION HANDLING  (Q&A prep, do not read aloud)

"Why pay when free apps exist?"
   They are ad-supported, generic, and have no live coaching. We are the only option
   here that puts a certified advisor on a call with you every month.

"Is my bank data safe?"
   Read-only sync, bank-grade encryption, no transaction authority. We can see the
   data; we can never move your money.

"What if I don't end up using it?"
   It is a monthly plan, so there is no lock-in. And the advisor call is a scheduled
   appointment, not an app you forget to open.

"4,999 upfront is a lot."
   That is exactly why the 499 monthly plan exists. The annual plan is for people who
   want the roughly 1,000-rupee saving.

"I already have a CA / advisor."
   They file and plan once a year. We handle the other 364 days: the daily spending,
   the alerts, the monthly course-correction.""",
]

# ================================================================= build
if __name__ == "__main__":
    for fn in (slide1, slide2, slide3, slide4, slide5):
        fn()
    for slide, note in zip(prs.slides, NOTES):
        slide.notes_slide.notes_text_frame.text = note
    prs.save(OUT)
    print("wrote", OUT)
